import traceback
from docx import Document
from pptx import Presentation
import PyPDF2  # Using PyPDF2 instead of PyMuPDF/fitz
from pdfminer.high_level import extract_text 

class TextExtractor:
    @staticmethod
    def extract_pdf(file_path: str) -> str:
        try:
            # Using PyPDF2 instead of fitz
            pdf_reader = PyPDF2.PdfReader(file_path)
            result = []
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                result.append(text)
            return "\n".join(result)
        except Exception as e:
            print(f"Error using PyPDF2: {e}")
            try:
                print("Falling back to pdfminer...")
                text = extract_text(file_path)
                return text
            except Exception as e:
                print(f"Error using pdfminer: {e}")
                return None

    @staticmethod
    def extract_docx(file_path: str) -> str:
        try:
            doc = Document(file_path)
            
            content = []
            
            for para in doc.paragraphs:
                if para.text.strip(): 
                    content.append({"text": para.text.strip()})

            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells] 
                    table_data.append(row_data)
                content.append({"table": table_data})

            return content

        except Exception as e:
            print(f"Error: {e}")
            return None

    @staticmethod
    def extract_pptx(file_path: str) -> str:
        try:
            presentation = Presentation(file_path)
            
            content = []
            for slide in presentation.slides:
                slide_content = []
                
                for shape in slide.shapes:
                    if shape.shape_type == 19: 
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells] 
                            table_data.append(row_data)
                        slide_content.append({"table": table_data})
                    elif hasattr(shape, "text") and shape.text.strip(): 
                        slide_content.append({"text": shape.text.strip()}) 
                    
                if slide_content:
                    content.append(slide_content)
            
            return content

        except Exception as e:
            print(f"Error: {e}")
            return None